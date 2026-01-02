import os
from langchain_chroma import Chroma
from langchain_community.document_loaders import PDFPlumberLoader
from langchain_community.document_loaders import TextLoader
from langchain_core.documents import Document as LcDocument
from langchain_community.document_loaders import UnstructuredHTMLLoader, UnstructuredEPubLoader, UnstructuredWordDocumentLoader
from pptx import Presentation
from app.services.service_utils import safe_detect_language

class DocumentLoadResult:
    def __init__(self, extension, size, doc_length, chunks_length, langs):
        self.extension = extension
        self.size = size
        self.doc_length = doc_length
        self.chunks_length = chunks_length
        self.lang = langs

    
class SessionDocumentLoadResult:
    def __init__(self, file, extension, size, doc_length, chunks_length, langs):
        self.file = file
        self.extension = extension
        self.size = size
        self.doc_length = doc_length
        self.chunks_length = chunks_length
        self.lang = langs

class DocumentLoader:
    def __init__(self, splitter, embedding, db_path):
        self.text_splitter = splitter
        self.embedding = embedding
        self.db_path = db_path

    def load_txt(self, path) -> DocumentLoadResult:
        txt = None
        with open(path, "r+") as f:
            txt = f.read()

        size_bytes = os.path.getsize(path)
        texts = self.text_splitter.split_text(txt)
        langs = safe_detect_language(chunks=texts)

        Chroma.from_texts(texts=texts, embedding=self.embedding, persist_directory=self.db_path)
        return DocumentLoadResult(extension='txt', size=size_bytes, doc_length=-1, chunks_length=len(texts), langs=langs)

    def load_txt_session(self, path, file, session_db_path) -> DocumentLoadResult:
        txt = None
        with open(path, "r+") as f:
            txt = f.read()

        size_bytes = os.path.getsize(path)
        texts = self.text_splitter.split_text(txt)
        langs = safe_detect_language(chunks=texts)

        Chroma.from_texts(texts=texts, embedding=self.embedding, persist_directory=session_db_path)
        return SessionDocumentLoadResult(file = file, extension='txt', size=size_bytes, doc_length=-1, chunks_length=len(texts), langs=langs)
    
    def load_pdf(self, path) -> DocumentLoadResult:
        size_bytes = os.path.getsize(path)
        pdf_loader = PDFPlumberLoader(path)
        docs = pdf_loader.load_and_split()
        chunks = self.text_splitter.split_documents(docs)
        langs = safe_detect_language(chunks=chunks)

        Chroma.from_documents(documents=chunks, embedding=self.embedding, persist_directory=self.db_path)
        return DocumentLoadResult(extension='pdf', size=size_bytes, doc_length=len(docs), chunks_length=len(chunks), langs=langs)

    def load_pdf_session(self, path, file, session_db_path) -> DocumentLoadResult:
        size_bytes = os.path.getsize(path)
        pdf_loader = PDFPlumberLoader(path)
        docs = pdf_loader.load_and_split()
        chunks = self.text_splitter.split_documents(docs)
        langs = safe_detect_language(chunks=chunks)

        Chroma.from_documents(documents=chunks, embedding=self.embedding, persist_directory=session_db_path)
        return SessionDocumentLoadResult(file=file, extension='pdf', size=size_bytes, doc_length=len(docs), chunks_length=len(chunks), langs=langs)

    def load_epub(self, path) -> DocumentLoadResult:
        size_bytes = os.path.getsize(path)
        epub_loader = UnstructuredEPubLoader(path)
        docs = epub_loader.load()
        chunks = self.text_splitter.split_documents(docs)
        langs = safe_detect_language(chunks=chunks)

        Chroma.from_documents(documents=chunks, embedding=self.embedding, persist_directory=self.db_path)
        return DocumentLoadResult(extension='epub', size=size_bytes, doc_length=len(docs), chunks_length=len(chunks), langs=langs)

    def load_epub_session(self, path, file, session_db_path) -> DocumentLoadResult:
        size_bytes = os.path.getsize(path)
        epub_loader = UnstructuredEPubLoader(path)
        docs = epub_loader.load()
        chunks = self.text_splitter.split_documents(docs)
        langs = safe_detect_language(chunks=chunks)

        Chroma.from_documents(documents=chunks, embedding=self.embedding, persist_directory=session_db_path)
        return SessionDocumentLoadResult(file=file, extension='epub', size=size_bytes, doc_length=len(docs), chunks_length=len(chunks), langs=langs)

    def load_docx(self, path) -> DocumentLoadResult:
        size_bytes = os.path.getsize(path)
        loader = UnstructuredWordDocumentLoader(path)
        docs = loader.load()
        chunks = self.text_splitter.split_documents(docs)
        langs = safe_detect_language(chunks=chunks)

        Chroma.from_documents(documents=chunks, embedding=self.embedding, persist_directory=self.db_path)
        return DocumentLoadResult(extension='docx', size=size_bytes, doc_length=1, chunks_length=len(chunks), langs=langs)
    
    def load_docx_session(self, path, file, session_db_path) -> DocumentLoadResult:
        size_bytes = os.path.getsize(path)
        loader = UnstructuredWordDocumentLoader(path)
        docs = loader.load()
        chunks = self.text_splitter.split_documents(docs)
        langs = safe_detect_language(chunks=chunks)

        Chroma.from_documents(documents=chunks, embedding=self.embedding, persist_directory=session_db_path)
        return SessionDocumentLoadResult(file=file, extension='docx', size=size_bytes, doc_length=1, chunks_length=len(chunks), langs=langs)

    def load_pptx(self, path) -> DocumentLoadResult:
        size_bytes = os.path.getsize(path)
        presentation = Presentation(path)
        documents = []

        for sld_num, sld in enumerate(presentation.slides, start=1):
            slide_text = []
            for shape in sld.shapes:
                if not shape.has_text_frame:
                    continue
                 
                for prg in shape.text_frame.paragraphs:
                    text = prg.text.strip()
                    if text:
                        slide_text.append(text)
            
            if slide_text:
                documents.append(
                    LcDocument( 
                        page_content="\n".join(slide_text),
                        metadata={
                            "slide": sld_num,
                            "type": "pptx"}
                ))

        chunks = self.text_splitter.split_documents(documents=documents)
        langs = safe_detect_language(chunks=chunks)

        Chroma.from_documents(documents=chunks, embedding=self.embedding, persist_directory=self.db_path)
        return DocumentLoadResult(extension='pptx', size=size_bytes, doc_length=len(presentation.slides), chunks_length=len(chunks), langs=langs)

    def load_pptx_session(self, path, file, session_db_path) -> DocumentLoadResult:
        size_bytes = os.path.getsize(path)
        presentation = Presentation(path)
        documents = []

        for sld_num, sld in enumerate(presentation.slides, start=1):
            slide_text = []
            for shape in sld.shapes:
                if not shape.has_text_frame:
                    continue
                 
                for prg in shape.text_frame.paragraphs:
                    text = prg.text.strip()
                    if text:
                        slide_text.append(text)
            
            if slide_text:
                documents.append(
                    LcDocument( 
                        page_content="\n".join(slide_text),
                        metadata={
                            "slide": sld_num,
                            "type": "pptx"}
                ))

        chunks = self.text_splitter.split_documents(documents=documents)
        langs = safe_detect_language(chunks=chunks)

        Chroma.from_documents(documents=chunks, embedding=self.embedding, persist_directory=session_db_path)
        return SessionDocumentLoadResult(file=file, extension='pptx', size=size_bytes, doc_length=len(presentation.slides), chunks_length=len(chunks), langs=langs)

    def load_md(self, path) -> DocumentLoadResult:
        size_bytes = os.path.getsize(path)
        loader = TextLoader(path, encoding="utf-8")
        docs = loader.load_and_split()
        chunks = self.text_splitter.split_documents(docs)
        langs = safe_detect_language(chunks=chunks)

        Chroma.from_documents(documents=chunks, embedding=self.embedding, persist_directory=self.db_path)        
        return DocumentLoadResult(extension='md', size=size_bytes, doc_length=len(docs), chunks_length=len(chunks), langs=langs)

    def load_md_session(self, path, file, session_db_path) -> DocumentLoadResult:
        size_bytes = os.path.getsize(path)
        loader = TextLoader(path, encoding="utf-8")
        docs = loader.load_and_split()
        chunks = self.text_splitter.split_documents(docs)
        langs = safe_detect_language(chunks=chunks)

        Chroma.from_documents(documents=chunks, embedding=self.embedding, persist_directory=session_db_path)
        return SessionDocumentLoadResult(file=file, extension='md', size=size_bytes, doc_length=len(docs), chunks_length=len(chunks), langs=langs)

    def load_html(self, path) -> DocumentLoadResult:
        size_bytes = os.path.getsize(path)
        loader = UnstructuredHTMLLoader(path)
        docs = loader.load()
        chunks = self.text_splitter.split_documents(documents=docs)
        langs = safe_detect_language(chunks=chunks)

        Chroma.from_documents(documents=chunks, embedding=self.embedding, persist_directory=self.db_path)
        return DocumentLoadResult(extension='html', size=size_bytes, doc_length=len(docs), chunks_length=len(chunks), langs=langs)
    
    def load_html_session(self, path, file, session_db_path) -> DocumentLoadResult:
        size_bytes = os.path.getsize(path)
        loader = UnstructuredHTMLLoader(path)
        docs = loader.load()
        chunks = self.text_splitter.split_documents(documents=docs)
        langs = safe_detect_language(chunks=chunks)

        Chroma.from_documents(documents=chunks, embedding=self.embedding, persist_directory=session_db_path)
        return SessionDocumentLoadResult(file=file, extension='html', size=size_bytes, doc_length=len(docs), chunks_length=len(chunks), langs=langs)